"""
RAG / knowledge-grounding layer.

parse -> clean -> chunk -> embed -> Chroma -> retrieve -> cited context block

Design notes
------------
* One Chroma collection per uploaded document (`doc_<id>`) so a session can be
  scoped to exactly the material the student uploaded.
* Embeddings run locally (MiniLM, CPU) — no embedding API key, works offline,
  and keeps the hackathon demo free.
* Every chunk keeps `page` and `section` metadata so the teacher can cite
  "page 41, section 4.2" instead of hallucinating.
"""
from __future__ import annotations

import logging
import re
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.config import settings

log = logging.getLogger("ai-teacher.rag")

SUPPORTED = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md"}

# A heading looks like "4.2 Ohm's Law", "CHAPTER 4", "Chapter 4 — Electricity"
_HEADING = re.compile(
    r"^\s*(?:(?:chapter|unit|adhyay|अध्याय)\s+\d+|\d+(?:\.\d+){0,2})\s*[.:\-–—]?\s+\S.{0,80}$",
    re.I,
)


# --------------------------------------------------------------------------- #
#  Data holders
# --------------------------------------------------------------------------- #
@dataclass
class DocMeta:
    doc_id: str
    filename: str
    pages: int
    chunks: int
    chars: int
    outline: list[str] = field(default_factory=list)
    preview: str = ""

    def as_dict(self) -> dict[str, Any]:
        return {
            "doc_id": self.doc_id,
            "filename": self.filename,
            "pages": self.pages,
            "chunks": self.chunks,
            "chars": self.chars,
            "outline": self.outline,
            "preview": self.preview,
        }


_DOCS: dict[str, DocMeta] = {}
_EMBEDDINGS: HuggingFaceEmbeddings | None = None


def get_embeddings() -> HuggingFaceEmbeddings:
    global _EMBEDDINGS
    if _EMBEDDINGS is None:
        log.info("loading embedding model %s", settings.embedding_model)
        _EMBEDDINGS = HuggingFaceEmbeddings(
            model_name=settings.embedding_model,
            encode_kwargs={"normalize_embeddings": True},
        )
    return _EMBEDDINGS


# --------------------------------------------------------------------------- #
#  1. Parsing
# --------------------------------------------------------------------------- #
def _clean(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    # de-hyphenate words broken across lines: "resis-\ntance" -> "resistance"
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    return text.strip()


def _parse_pdf(path: Path) -> list[tuple[int, str]]:
    import fitz  # PyMuPDF

    out: list[tuple[int, str]] = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, start=1):
            out.append((i, _clean(page.get_text("text"))))
    return out


def _parse_docx(path: Path) -> list[tuple[int, str]]:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        style = (para.style.name or "").lower()
        parts.append(f"\n{para.text.strip()}\n" if "heading" in style else para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return [(1, _clean("\n".join(parts)))]


def _parse_pptx(path: Path) -> list[tuple[int, str]]:
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text.strip())
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            chunks.append("[notes] " + slide.notes_slide.notes_text_frame.text.strip())
        out.append((i, _clean("\n".join(chunks))))
    return out


def _parse_text(path: Path) -> list[tuple[int, str]]:
    return [(1, _clean(path.read_text(encoding="utf-8", errors="ignore")))]


PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".doc": _parse_docx,
    ".pptx": _parse_pptx,
    ".ppt": _parse_pptx,
    ".txt": _parse_text,
    ".md": _parse_text,
}


def extract_pages(path: Path) -> list[tuple[int, str]]:
    ext = path.suffix.lower()
    if ext not in PARSERS:
        raise ValueError(f"unsupported file type '{ext}'. Supported: {sorted(SUPPORTED)}")
    return [(p, t) for p, t in PARSERS[ext](path) if t.strip()]


def extract_outline(pages: list[tuple[int, str]], limit: int = 40) -> list[str]:
    """Cheap structural map of the document — feeds the lesson planner so it
    can say 'Chapter 4' and mean it."""
    seen: list[str] = []
    for page_no, text in pages:
        for line in text.splitlines():
            line = line.strip()
            if 4 < len(line) < 90 and _HEADING.match(line):
                entry = f"p{page_no}: {line}"
                if entry not in seen:
                    seen.append(entry)
            if len(seen) >= limit:
                return seen
    return seen


# --------------------------------------------------------------------------- #
#  2. Chunk + embed + store
# --------------------------------------------------------------------------- #
def _collection(doc_id: str) -> Chroma:
    return Chroma(
        collection_name=f"doc_{doc_id}",
        embedding_function=get_embeddings(),
        persist_directory=str(settings.chroma_path),
    )


def ingest(path: Path, original_name: str) -> DocMeta:
    doc_id = uuid.uuid4().hex[:12]
    pages = extract_pages(path)
    if not pages:
        raise ValueError("no extractable text found (is this a scanned PDF? OCR is not enabled)")

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ". ", "। ", " ", ""],
    )

    docs: list[Document] = []
    current_section = ""
    for page_no, text in pages:
        for line in text.splitlines():
            if _HEADING.match(line.strip()):
                current_section = line.strip()
                break
        for piece in splitter.split_text(text):
            docs.append(
                Document(
                    page_content=piece,
                    metadata={
                        "doc_id": doc_id,
                        "filename": original_name,
                        "page": page_no,
                        "section": current_section,
                    },
                )
            )

    store = _collection(doc_id)
    # batch to stay under Chroma's max add size on big textbooks
    for i in range(0, len(docs), 128):
        store.add_documents(docs[i : i + 128])

    meta = DocMeta(
        doc_id=doc_id,
        filename=original_name,
        pages=len(pages),
        chunks=len(docs),
        chars=sum(len(t) for _, t in pages),
        outline=extract_outline(pages),
        preview=pages[0][1][:600],
    )
    _DOCS[doc_id] = meta
    log.info("ingested %s -> %s (%s chunks)", original_name, doc_id, len(docs))
    return meta


def get_doc(doc_id: str) -> DocMeta | None:
    return _DOCS.get(doc_id)


def list_docs() -> list[dict[str, Any]]:
    return [m.as_dict() for m in _DOCS.values()]


# --------------------------------------------------------------------------- #
#  3. Retrieval
# --------------------------------------------------------------------------- #
def retrieve(doc_id: str, query: str, k: int | None = None) -> list[Document]:
    if not doc_id:
        return []
    k = k or settings.retrieve_k
    try:
        return _collection(doc_id).similarity_search(query, k=k)
    except Exception as exc:  # noqa: BLE001
        log.warning("retrieval failed for %s: %s", doc_id, exc)
        return []


def build_context(docs: list[Document], max_chars: int = 6000) -> tuple[str, list[dict[str, Any]]]:
    """Returns (prompt_block, citations)."""
    if not docs:
        return "", []
    lines: list[str] = ["SOURCE MATERIAL (ground every claim in this):"]
    citations: list[dict[str, Any]] = []
    used = 0
    for i, d in enumerate(docs, start=1):
        page = d.metadata.get("page", "?")
        section = d.metadata.get("section", "")
        header = f"[{i}] page {page}{(' | ' + section) if section else ''}"
        body = d.page_content.strip()
        if used + len(body) > max_chars:
            body = body[: max(0, max_chars - used)]
        if not body:
            break
        lines.append(f"{header}\n{body}")
        citations.append(
            {
                "n": i,
                "page": page,
                "section": section,
                "filename": d.metadata.get("filename", ""),
                "snippet": body[:220],
            }
        )
        used += len(body)
        if used >= max_chars:
            break
    lines.append(
        "END OF SOURCE MATERIAL. If something you need is missing from the "
        "material above, say so in half a sentence and then teach it from "
        "first principles — never invent a citation."
    )
    return "\n\n".join(lines), citations


def context_for(doc_id: str, query: str, k: int | None = None) -> tuple[str, list[dict[str, Any]]]:
    return build_context(retrieve(doc_id, query, k))
